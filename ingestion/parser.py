import uuid
from pathlib import Path
from config import SUPPORTED_EXTENSIONS


def extract(file_path: str) -> dict:
    path = Path(file_path)
    ext  = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}")

    parsers = {
        ".pdf":  _parse_pdf,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".csv":  _parse_csv,
        ".txt":  _parse_txt,
    }

    blocks = parsers[ext](path)

    return {
        "doc_id":    str(uuid.uuid4()),
        "filename":  path.name,
        "ext":       ext,
        "blocks":    blocks,
        "full_text": _build_full_text(blocks),
    }


# ── PDF ───────────────────────────────────────────────────────────────────────

def _parse_pdf(path: Path) -> list[dict]:
    import fitz

    doc = fitz.open(str(path))
    blocks = []
    char_cursor = 0

    for page_num, page in enumerate(doc, start=1):
        page_dict = page.get_text("dict")

        for block in page_dict["blocks"]:
            if block["type"] != 0:
                continue

            lines = []
            for line in block["lines"]:
                line_text = " ".join(
                    span["text"] for span in line["spans"]
                ).strip()
                if line_text:
                    lines.append(line_text)

            text = "\n".join(lines).strip()
            if not text:
                continue

            first_size = block["lines"][0]["spans"][0]["size"] \
                if block["lines"] else 0
            elem_type  = "heading" if first_size >= 14 else "paragraph"

            start = char_cursor
            end   = char_cursor + len(text)
            char_cursor = end + 1

            blocks.append({
                "type":       elem_type,
                "text":       text,
                "page":       page_num,
                "start_char": start,
                "end_char":   end,
            })

    doc.close()
    return blocks


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _parse_docx(path: Path) -> list[dict]:
    from docx import Document

    doc = Document(str(path))
    blocks = []
    char_cursor = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style     = para.style.name
        elem_type = "heading" if "heading" in style.lower() else "paragraph"

        start = char_cursor
        end   = char_cursor + len(text)
        char_cursor = end + 1

        blocks.append({
            "type":       elem_type,
            "text":       text,
            "page":       None,
            "style":      style,
            "start_char": start,
            "end_char":   end,
        })

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        text = "\n".join(rows).strip()
        if not text:
            continue

        start = char_cursor
        end   = char_cursor + len(text)
        char_cursor = end + 1

        blocks.append({
            "type":       "table",
            "text":       text,
            "page":       None,
            "style":      "Table",
            "start_char": start,
            "end_char":   end,
        })

    return blocks


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _parse_pdf(path: Path) -> list[dict]:
    import fitz

    doc = fitz.open(str(path))
    blocks = []
    char_cursor = 0

    for page_num, page in enumerate(doc, start=1):
        page_dict = page.get_text("dict")

        for block in page_dict["blocks"]:
            if block["type"] != 0:
                continue

            lines = []
            for line in block["lines"]:
                line_text = " ".join(
                    span["text"] for span in line["spans"]
                ).strip()
                if line_text:
                    lines.append(line_text)

            text = "\n".join(lines).strip()
            if not text:
                continue

            first_size = block["lines"][0]["spans"][0]["size"] \
                if block["lines"] else 0
            elem_type  = "heading" if first_size >= 14 else "paragraph"

            start = char_cursor
            end   = char_cursor + len(text)
            char_cursor = end + 1

            blocks.append({
                "type":       elem_type,
                "text":       text,
                "page":       page_num,
                "start_char": start,
                "end_char":   end,
            })

    doc.close()

    # guard: if almost no text was extracted, PDF is likely image-based
    full_text = _build_full_text(blocks)
    if len(full_text.strip()) < 50:
        raise ValueError(
            f"'{path.name}' appears to be an image-based or scanned PDF with no "
            f"extractable text. Please use a PDF with selectable text, or run OCR "
            f"on it first before uploading."
        )

    return blocks
# ── PPTX ──────────────────────────────────────────────────────────────────────

def _parse_pptx(path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks = []
    char_cursor = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            try:
                is_title = (
                    shape.is_placeholder
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
            except Exception:
                is_title = False

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                elem_type   = "heading" if is_title else "paragraph"
                start       = char_cursor
                end         = char_cursor + len(text)
                char_cursor = end + 1

                blocks.append({
                    "type":       elem_type,
                    "text":       text,
                    "page":       slide_num,
                    "start_char": start,
                    "end_char":   end,
                })


    return blocks

# ── XLSX ──────────────────────────────────────────────────────────────────────


_XLSX_ROWS_PER_CHUNK = 80


def _parse_xlsx(path: Path) -> list[dict]:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=True)
    blocks = []
    char_cursor = 0

    for sheet in wb.worksheets:
        all_rows = []
        for row in sheet.iter_rows(values_only=True):
            cells    = [str(c).strip() if c is not None else "" for c in row]
            row_text = " | ".join(cells).strip(" |")
            if row_text:
                all_rows.append(row_text)

        if not all_rows:
            continue

        header    = all_rows[0]
        data_rows = all_rows[1:]

        # small sheet — single chunk
        full_text = f"[Sheet: {sheet.title}]\n" + "\n".join(all_rows)
        if len(data_rows) <= 50:
            start = char_cursor
            end   = char_cursor + len(full_text)
            char_cursor = end + 1
            blocks.append({
                "type":       "table",
                "text":       full_text,
                "page":       sheet.title,
                "start_char": start,
                "end_char":   end,
            })
            continue

        # large sheet — chunk by row count, header repeated every chunk
        for i in range(0, len(data_rows), _XLSX_ROWS_PER_CHUNK):
            row_group = data_rows[i: i + _XLSX_ROWS_PER_CHUNK]
            text      = (
                f"[Sheet: {sheet.title}  "
                f"rows {i+1}-{i+len(row_group)}]\n"
                f"{header}\n"
                + "\n".join(row_group)
            )
            start = char_cursor
            end   = char_cursor + len(text)
            char_cursor = end + 1
            blocks.append({
                "type":       "table",
                "text":       text,
                "page":       sheet.title,
                "start_char": start,
                "end_char":   end,
            })

    return blocks


# ── CSV ───────────────────────────────────────────────────────────────────────

# Same row-per-chunk constant as XLSX
_CSV_ROWS_PER_CHUNK = 80


def _parse_csv(path: Path) -> list[dict]:
    import csv

    blocks = []
    char_cursor = 0

    with open(str(path), newline="", encoding="utf-8-sig") as f:
        reader   = csv.reader(f)
        all_rows = []
        for row in reader:
            row_text = " | ".join(cell.strip() for cell in row)
            if row_text.strip(" |"):
                all_rows.append(row_text)

    if not all_rows:
        return blocks

    header    = all_rows[0]
    data_rows = all_rows[1:]

    # small CSV — single chunk
    if len(data_rows) <= 50:
        full_text = f"[CSV: {path.name}]\n" + "\n".join(all_rows)
        start     = char_cursor
        end       = char_cursor + len(full_text)
        blocks.append({
            "type":       "table",
            "text":       full_text,
            "page":       path.name,
            "start_char": start,
            "end_char":   end,
        })
        return blocks

    # large CSV — chunk by row count, header repeated every chunk
    for i in range(0, len(data_rows), _CSV_ROWS_PER_CHUNK):
        row_group = data_rows[i: i + _CSV_ROWS_PER_CHUNK]
        text      = (
            f"[CSV: {path.name}  "
            f"rows {i+1}-{i+len(row_group)}]\n"
            f"{header}\n"
            + "\n".join(row_group)
        )
        start = char_cursor
        end   = char_cursor + len(text)
        char_cursor = end + 1
        blocks.append({
            "type":       "table",
            "text":       text,
            "page":       path.name,
            "start_char": start,
            "end_char":   end,
        })

    return blocks


# ── TXT ───────────────────────────────────────────────────────────────────────

def _parse_txt(path: Path) -> list[dict]:
    blocks = []
    char_cursor = 0

    with open(str(path), encoding="utf-8", errors="ignore") as f:
        raw = f.read()

    paragraphs = [p.strip() for p in raw.split("\n\n") if p.strip()]

    for para in paragraphs:
        lines     = para.splitlines()
        elem_type = (
            "heading"
            if len(lines) == 1 and len(para) < 80 and not para.endswith(".")
            else "paragraph"
        )

        start = char_cursor
        end   = char_cursor + len(para)
        char_cursor = end + 1

        blocks.append({
            "type":       elem_type,
            "text":       para,
            "page":       None,
            "start_char": start,
            "end_char":   end,
        })

    return blocks


# ── HELPERS ───────────────────────────────────────────────────────────────────

def _build_full_text(blocks: list[dict]) -> str:
    return "\n".join(b["text"] for b in blocks)